"""Office document extractors (docx / xlsx / pptx).

Optional extra (``pleno-pii-scanner[office]``). Each format gets its own
class so callers can selectively register what they want — e.g. a
SharePoint connector that only ingests Word docs doesn't need to import
openpyxl.

Path hints:
- docx -> "paragraph:N"
- xlsx -> "sheet:NAME!cell:A1" — preserves cell coordinates so ops can
  jump directly to the offending cell in Excel
- pptx -> "slide:N/shape:M"
"""

from __future__ import annotations

from collections.abc import AsyncIterator

from pleno_pii_scanner.extractors.base import (
    ExtractedFragment,
    ExtractorError,
    doc_payload,
)
from pleno_pii_scanner.sources.base import Document, DocumentChunk


class DocxExtractor:
    """Word .docx -> per-paragraph fragments."""

    name = "office:docx"
    accepts = frozenset(
        {
            "application/vnd.openxmlformats-officedocument."
            "wordprocessingml.document",
        }
    )

    def __init__(self) -> None:
        try:
            import docx
        except ImportError as exc:
            raise ExtractorError(
                "DocxExtractor requires the [office] extra"
            ) from exc
        self._docx = docx

    async def extract(
        self, doc: Document | DocumentChunk
    ) -> AsyncIterator[ExtractedFragment]:
        payload = _as_bytes(doc)
        import io

        document = self._docx.Document(io.BytesIO(payload))
        for idx, paragraph in enumerate(document.paragraphs):
            text = paragraph.text
            if not text:
                continue
            yield ExtractedFragment(
                text=text,
                path_hint=f"paragraph:{idx + 1}",
                byte_offset=None,
                extractor=self.name,
            )


class XlsxExtractor:
    """Excel .xlsx -> per-cell fragments (string cells only).

    Numeric / boolean / date cells cannot carry PII (they pass through
    Excel's coercion to a strict type), so we skip them. This matches
    the columnar extractor's "string columns only" rule from ADR-0007 §6.
    """

    name = "office:xlsx"
    accepts = frozenset(
        {
            "application/vnd.openxmlformats-officedocument."
            "spreadsheetml.sheet",
        }
    )

    def __init__(self) -> None:
        try:
            import openpyxl
        except ImportError as exc:
            raise ExtractorError(
                "XlsxExtractor requires the [office] extra"
            ) from exc
        self._openpyxl = openpyxl

    async def extract(
        self, doc: Document | DocumentChunk
    ) -> AsyncIterator[ExtractedFragment]:
        import io

        payload = _as_bytes(doc)
        wb = self._openpyxl.load_workbook(
            io.BytesIO(payload), read_only=True, data_only=True
        )
        try:
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                for row in ws.iter_rows(values_only=False):
                    for cell in row:
                        if cell.value is None or not isinstance(
                            cell.value, str
                        ):
                            continue
                        yield ExtractedFragment(
                            text=cell.value,
                            path_hint=f"sheet:{sheet_name}!cell:{cell.coordinate}",
                            byte_offset=None,
                            extractor=self.name,
                        )
        finally:
            wb.close()


class PptxExtractor:
    """PowerPoint .pptx -> per-shape text fragments."""

    name = "office:pptx"
    accepts = frozenset(
        {
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation",
        }
    )

    def __init__(self) -> None:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ExtractorError(
                "PptxExtractor requires the [office] extra"
            ) from exc
        self._Presentation = Presentation

    async def extract(
        self, doc: Document | DocumentChunk
    ) -> AsyncIterator[ExtractedFragment]:
        import io

        payload = _as_bytes(doc)
        prs = self._Presentation(io.BytesIO(payload))
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape_idx, shape in enumerate(slide.shapes, start=1):
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text
                if not text:
                    continue
                yield ExtractedFragment(
                    text=text,
                    path_hint=f"slide:{slide_idx}/shape:{shape_idx}",
                    byte_offset=None,
                    extractor=self.name,
                )


def _as_bytes(doc: Document | DocumentChunk) -> bytes:
    payload = doc_payload(doc)
    if isinstance(payload, str):
        raise ExtractorError("Office extractors require binary payload")
    return payload
